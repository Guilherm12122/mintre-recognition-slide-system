from pptx import Presentation
import glob

def extract_data_from_presentation_slide(slide_path: str):
    
    apresentacao_slide = Presentation(slide_path)

    objeto_dicionario = {"path_slide": slide_path}
    lista_dados_slide = []

    for index, slide in enumerate(apresentacao_slide.slides):

        texto_extraido = ''

        for objeto in slide.shapes:

            if hasattr(objeto, 'text'):

                texto_extraido += objeto.text

        lista_dados_slide.append({ "posicao": index + 1, "conteudo": texto_extraido })

    objeto_dicionario['data'] = lista_dados_slide

    return objeto_dicionario

# Obtêm o caminho de todos os slides de uma pasta.
def get_data_from_files_pptx(general_path: str):

    return [extract_data_from_presentation_slide(file_path) for file_path in glob.glob(general_path, recursive=True)]



